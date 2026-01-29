from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw
import os

# --- 1. Helper: Generate Tech Background Image ---
def create_tech_background(filename):
    width = 1920
    height = 1080
    # Dark blue/black background
    img = Image.new('RGB', (width, height), color=(10, 15, 30))
    draw = ImageDraw.Draw(img)
    
    # Draw faint grid
    grid_color = (30, 40, 60)
    step = 50
    for x in range(0, width, step):
        draw.line((x, 0, x, height), fill=grid_color)
    for y in range(0, height, step):
        draw.line((0, y, width, y), fill=grid_color)
        
    # Draw some random tech lines/accents
    accent_color = (0, 100, 200)
    draw.line((0, 100, width, 100), fill=accent_color, width=2)
    draw.line((0, height-100, width, height-100), fill=accent_color, width=2)
    
    img.save(filename)
    return filename

# --- 2. Helper: Slide Creation ---
def set_slide_background(slide, image_path):
    left = top = 0
    pic = slide.shapes.add_picture(image_path, left, top, width=prs.slide_width, height=prs.slide_height)
    # Move to back (this is a hack, usually adding first puts it at back, but we add shapes later)
    # python-pptx doesn't strictly support z-ordering easily, but adding first usually works.
    
    # Move other shapes to front is hard, so we just add bg first.

def create_slide(prs, layout_index, title_text, bg_image):
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    
    # Add background image FIRST so it's behind everything
    set_slide_background(slide, bg_image)
    
    # Re-add title and content placeholders if they are covered or need styling
    # Actually, adding an image covers the placeholders. 
    # Better approach: Add image, then add text boxes manually or try to use master slide.
    # For simplicity in python-pptx without master slide manipulation:
    # We will add text boxes manually to ensure visibility and z-order.
    
    return slide

def add_title(slide, text):
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 255, 255) # Cyan
    p.alignment = PP_ALIGN.LEFT

def add_subtitle(slide, text):
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200) # Light Gray

def add_content_box(slide, items):
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = 'Calibri'
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255) # White
        p.space_after = Pt(14)

# --- 3. Main Script ---

# Generate background
bg_filename = "tech_bg.png"
create_tech_background(bg_filename)

prs = Presentation()
# Set 16:9 aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1: Cover
slide = create_slide(prs, 6, "", bg_filename) # 6 is usually blank
add_title(slide, "SandAgent")
# Center the title for cover
slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
slide.shapes[-1].top = Inches(2.5)
slide.shapes[-1].left = Inches(2) # Center roughly

# Subtitle
left = Inches(2)
top = Inches(4)
width = Inches(9.333)
height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Turn powerful coding agents into universal Super Agents"
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(100, 200, 255)
p.alignment = PP_ALIGN.CENTER

# Slide 2: The Problem
slide = create_slide(prs, 6, "", bg_filename)
add_title(slide, "The Problem: Building Agents is Hard")
add_content_box(slide, [
    "Traditional SDK approach is painful and slow (6+ months)",
    "Complex Memory Management (History, Context, Summarization)",
    "Difficult Tool Integration & MCP Servers",
    "Infrastructure Headaches (Sandboxes, Filesystems)",
    "Endless Prompt Engineering"
])

# Slide 3: The Solution
slide = create_slide(prs, 6, "", bg_filename)
add_title(slide, "The Solution: SandAgent")
add_content_box(slide, [
    "Don't rebuild the agent — just redirect it.",
    "Reuse Coding Agents: Leverage Claude Code, Codex CLI",
    "Skip the hard parts: Memory, Context, Tools are built-in",
    "Speed: 1 day to production vs 6 months",
    "Simplicity: Define agents with simple Markdown templates"
])

# Slide 4: Use Cases
slide = create_slide(prs, 6, "", bg_filename)
add_title(slide, "Specialized Super Agents")
add_content_box(slide, [
    "Data Analyst: SQL, Python, Visualization ('analyst')",
    "Research Assistant: Web research, Source evaluation ('researcher')",
    "Code Assistant: Review, Debug, Refactor ('coder')",
    "SEO Agent: Keyword research, Optimization ('seo-agent')",
    "Custom: Define any role via CLAUDE.md"
])

# Slide 5: Architecture (Visual)
slide = create_slide(prs, 6, "", bg_filename)
add_title(slide, "How It Works")

# Draw simple flow
# Shape 1: Template
left = Inches(1)
top = Inches(3)
width = Inches(2.5)
height = Inches(1.5)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
shape.text = "Template\n(CLAUDE.md)"
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 100, 200)
shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Arrow 1
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.6), Inches(3.6), Inches(1), Inches(0.3))

# Shape 2: Runner
left = Inches(4.8)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
shape.text = "Runner\n(Claude/Codex)"
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 150, 100)

# Arrow 2
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.4), Inches(3.6), Inches(1), Inches(0.3))

# Shape 3: Sandbox
left = Inches(8.6)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
shape.text = "Sandbox\n(E2B / Sandock)"
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(150, 50, 150)

# Slide 6: Key Features
slide = create_slide(prs, 6, "", bg_filename)
add_title(slide, "Key Features")
add_content_box(slide, [
    "Template-Based: Zero code, just Markdown",
    "Sandboxed Execution: Secure, isolated environments",
    "Persistent Sessions: Resume work anytime",
    "Swappable Infrastructure: Local / Docker / Cloud",
    "GAIA Benchmark Included"
])

# Slide 7: Quick Start
slide = create_slide(prs, 6, "", bg_filename)
add_title(slide, "Get Started")
add_content_box(slide, [
    "1. Clone Repo & Install Dependencies",
    "2. Configure API Keys (Anthropic)",
    "3. Run Web UI: pnpm dev -> localhost:3000",
    "Or use CLI: sandagent run -- 'Make a PPT'"
])

# Slide 8: End
slide = create_slide(prs, 6, "", bg_filename)
add_title(slide, "SandAgent")
slide.shapes[-1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
slide.shapes[-1].top = Inches(3)
slide.shapes[-1].left = Inches(2)

left = Inches(2)
top = Inches(4.5)
width = Inches(9.333)
height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "https://github.com/vikadata/sandagent"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(100, 200, 255)
p.alignment = PP_ALIGN.CENTER

# Save
try:
    prs.save('SandAgent_Presentation_Tech.pptx')
    print("Presentation saved successfully as SandAgent_Presentation_Tech.pptx")
except Exception as e:
    print(f"Error saving presentation: {e}")

# Cleanup
if os.path.exists(bg_filename):
    os.remove(bg_filename)
